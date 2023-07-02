import os
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
import os
# ベースフォルダのパスを入力させる
base_folder = input("ベースフォルダのパスを入力してください: ")
# ベースフォルダ内のサブフォルダを取得
image_folders = [f1 for f1 in os.listdir(base_folder) if os.path.isdir(os.path.join(base_folder, f1))]
template_file = 'test_template.pptx'
output_file = 'output.pptx'
# パワーポイントプレゼンテーションを作成
presentation = Presentation(template_file)
# 画像フォルダ内の画像ファイルを取得
for image_folder in image_folders:
    image_files = [f for f in os.listdir(os.path.join(base_folder, image_folder)) if f.endswith(('.jpg', '.jpeg', '.png', '.gif'))]
    slide_layout = presentation.slide_layouts[6]  # レイアウト6を使用するスライドレイアウト
    slide = presentation.slides.add_slide(slide_layout)
    # スライド上に画像を配置する位置とサイズの設定
    left = Cm(1)  # 左端の位置（センチメートル単位）
    top = Cm(2)  # 上端の位置（センチメートル単位）
    width = Cm(5)  # 画像の幅（センチメートル単位）
    height = Cm(4)  # 画像の高さ（センチメートル単位）
    space = Cm(1)  # 画像間のスペース（センチメートル単位
    # フォルダ名を記載するテキストボックスを作成
    folder_name_textbox_left = Cm(12)
    folder_name_textbox_top = Cm(6)
    folder_name_textbox_width = Cm(8)
    folder_name_textbox_height = Cm(1)
    folder_name_textbox = slide.shapes.add_textbox(folder_name_textbox_left, folder_name_textbox_top, folder_name_textbox_width, folder_name_textbox_height)
    folder_name_text_frame = folder_name_textbox.text_frame
    folder_name_text_frame.word_wrap = False
    folder_name_text = folder_name_text_frame.add_paragraph().add_run()
    folder_name_text.text = image_folder
    folder_name_text.font.size = Pt(28)
    # README.mdのテキストを読み込む
    image_folder=os.path.join(base_folder, image_folder)
    readme_file = os.path.join(image_folder, 'README.md')
    if os.path.isfile(readme_file):
        with open(readme_file, 'r', encoding='utf-8') as f:
            readme_text = f.read()
    else:
        readme_text = ""
    # 画像をスライドに配置する
    for i, image_file in enumerate(image_files):
        if i % 9 == 0:
            slide_layout = presentation.slide_layouts[6]  # レイアウト6を使用するスライドレイアウト
            slide = presentation.slides.add_slide(slide_layout)
            left = Cm(1)  # 左端の位置（センチメートル単位）
            top = Cm(2)  # 上端の位置（センチメートル単位）
            width = Cm(5)  # 画像の幅（センチメートル単位）
            height = Cm(4)  # 画像の高さ（センチメートル単位）
            space = Cm(1)  # 画像間のスペース（センチメートル単位
            folder_name_textbox_left = Cm(1)
            folder_name_textbox_top = Cm(0.1)
            folder_name_textbox_width = Cm(8)
            folder_name_textbox_height = Cm(1)
            folder_name_textbox = slide.shapes.add_textbox(folder_name_textbox_left, folder_name_textbox_top, folder_name_textbox_width, folder_name_textbox_height)
            folder_name_text_frame = folder_name_textbox.text_frame
            folder_name_text_frame.word_wrap = False
            folder_name_text = folder_name_text_frame.add_paragraph().add_run()
            folder_name_text.text = image_folder
            folder_name_text.font.size = Pt(18)
            # テキストボックスを作成してテキストを貼り付ける
            textbox_left = Cm(19)  # 左端の位置（センチメートル単位）
            textbox_top = Cm(1)  # 上端の位置（センチメートル単位）
            textbox_width = presentation.slide_width - textbox_left - Cm(1)  # テキストボックスの幅
            textbox_height = presentation.slide_height - textbox_top - Cm(1)  # テキストボックスの高さ
            textbox = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
            # textbox = presentation.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
            textbox.text = readme_text
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            # テキストボックスのテキストを設定
            p = text_frame.add_paragraph()
            p.text = readme_text
            # テキストの書式設定
            p.font.size = Pt(12)  # フォントサイズを14ポイントに設定
            p.font.name = "Meiryo"  # フォント名をMeiryoに設定
        image_path = os.path.join(image_folder, image_file)
        image_path= os.path.join(base_folder, image_path)
        image = slide.shapes.add_picture(image_path, left, top, width, height)
        left += width + space  # 次の画像の左端位置を更新
        k=(i+1) % 3
        if k == 0:
            left=Cm(1)
            top += height + space  # 次の行の上端位置を更新
presentation.save(output_file)
os.startfile(output_file)